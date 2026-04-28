"""A2A Handler for FileClaw - Universal File Import/Export for All Agents"""
import sys, json, csv, zipfile
from pathlib import Path
from datetime import datetime

PROJECT_ROOT = Path(__file__).resolve().parent.parent.parent
EXPORTS = PROJECT_ROOT / "exports"
sys.path.insert(0, str(PROJECT_ROOT))
from shared.base_agent import BaseAgent
from shared.security import InputSanitizer

class FileClawAgent(BaseAgent):
    def __init__(self):
        super().__init__("fileclaw")
        self.text_formats = {".md",".txt",".html",".xml",".json",".csv",".yaml",".yml",".toml",".ini",".cfg",".py",".js",".ts",".rs",".go",".cpp",".c",".java",".rb",".php",".css",".sql",".r",".sh",".bat",".ps1"}
        self.binary_importers = {
            ".pdf":self._read_pdf, ".docx":self._read_docx, ".xlsx":self._read_xlsx,
            ".pptx":self._read_pptx, ".rtf":self._read_rtf,
            ".png":self._read_image, ".jpg":self._read_image, ".jpeg":self._read_image,
            ".gif":self._read_image, ".bmp":self._read_image, ".webp":self._read_image,
            ".svg":self._read_image, ".tiff":self._read_image,
            ".zip":self._read_archive, ".tar":self._read_archive, ".gz":self._read_archive,
            ".mp3":self._read_media, ".mp4":self._read_media, ".wav":self._read_media,
            ".mkv":self._read_media, ".avi":self._read_media, ".mov":self._read_media,
            ".mpeg":self._read_media, ".mpg":self._read_media
        }

    def _gather_context(self, query=""):
        parts = []
        web = self.call_agent("webclaw", f"search file format {query}", timeout=15)
        if web: parts.append("[WebClaw]: " + web)
        data = self.call_agent("dataclaw", f"search {query}", timeout=15)
        if data: parts.append("[DataClaw]: " + data)
                # Search chronicle index
        chronicle_results = self.search_chronicle(query, limit=2000000)
        if chronicle_results:
            for c in chronicle_results:
                if hasattr(c, "url"):
                    parts.append(c.url)

        return "\n".join(parts)

    def _format_size(self, s):
        for u in ["B","KB","MB","GB"]:
            if s < 1024: return f"{s:.1f}{u}"
            s /= 1024
        return f"{s:.1f}TB"

    def _read_pdf(self, p):
        try:
            from pypdf import PdfReader
            r = PdfReader(str(p)); t = f"PDF: {p.name}, {len(r.pages)} pages\n\n"
            for i, pg in enumerate(r.pages): t += f"--- Page {i+1} ---\n{pg.extract_text() or ''}\n"
            return t
        except Exception as e: return f"PDF error: {e}"

    def _read_docx(self, p):
        try:
            from docx import Document
            d = Document(str(p)); t = f"DOCX: {p.name}\n\n"
            for para in d.paragraphs: t += para.text + "\n"
            return t
        except Exception as e: return f"DOCX error: {e}"

    def _read_xlsx(self, p):
        try:
            import pandas as pd
            sheets = pd.read_excel(p, sheet_name=None)
            t = f"XLSX: {p.name}, {len(sheets)} sheets\n\n"
            for n, df in sheets.items(): t += f"--- {n} ---\n{df.head(20).to_string()}\n\n"
            return t
        except Exception as e: return f"XLSX error: {e}"

    def _read_pptx(self, p):
        try:
            from pptx import Presentation
            prs = Presentation(str(p)); t = f"PPTX: {p.name}, {len(prs.slides)} slides\n\n"
            for i, sl in enumerate(prs.slides):
                t += f"--- Slide {i+1} ---\n"
                for sh in sl.shapes:
                    if hasattr(sh, "text") and sh.text: t += sh.text + "\n"
            return t
        except Exception as e: return f"PPTX error: {e}"

    def _read_rtf(self, p):
        try:
            import re
            c = p.read_text(encoding="utf-8", errors="replace")
            c = re.sub(r"\\[a-z]+\d*", "", c)
            return f"RTF: {p.name}\n\n{c}"
        except Exception as e: return f"RTF error: {e}"

    def _read_image(self, p):
        try:
            from PIL import Image
            img = Image.open(p)
            info = f"IMAGE: {p.name}\n  Format: {img.format}\n  Size: {img.size[0]}x{img.size[1]}\n  Mode: {img.mode}\n  File: {self._format_size(p.stat().st_size)}"
            try:
                exif = img._getexif()
                if exif: info += f"\n  EXIF: {len(exif)} tags"
            except: pass
            return info
        except Exception as e: return f"Image error: {e}"

    def _read_archive(self, p):
        try:
            t = f"ARCHIVE: {p.name}\n  Size: {self._format_size(p.stat().st_size)}\n\nContents:\n"
            if p.suffix == ".zip":
                with zipfile.ZipFile(p) as zf:
                    for info in zf.infolist():
                        t += f"  {'[D]' if info.is_dir() else '[F]'} {info.filename} ({self._format_size(info.file_size)})\n"
            else:
                import tarfile
                with tarfile.open(p) as tf:
                    for m in tf.getmembers():
                        t += f"  {'[D]' if m.isdir() else '[F]'} {m.name} ({self._format_size(m.size)})\n"
            return t
        except Exception as e: return f"Archive error: {e}"

    def _read_media(self, p):
        return f"MEDIA: {p.name}\n  Type: {p.suffix}\n  Size: {self._format_size(p.stat().st_size)}"

    def _import(self, fp):
        p = Path(fp)
        if not p.exists(): return f"Error: File not found: {fp}"
        ext = p.suffix.lower()
        if ext in self.binary_importers: return self.binary_importers[ext](p)
        if ext in self.text_formats:
            try:
                c = p.read_text(encoding="utf-8", errors="replace")
                return f"FILE: {p.name}\n  Size: {self._format_size(p.stat().st_size)}\n\n{c}"
            except: return f"Read error: {p.name}"
        try: return f"UNKNOWN: {p.name}\n{p.read_text(encoding='utf-8', errors='replace')}"
        except: return f"UNKNOWN BINARY: {p.name}\n  Type: {ext}\n  Size: {self._format_size(p.stat().st_size)}"

    def _export(self, fmt, content, name=""):
        EXPORTS.mkdir(exist_ok=True)
        fmt = fmt.lower().lstrip(".")
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        fn = f"{name or 'export'}_{ts}.{fmt}"
        fn = InputSanitizer.sanitize_filename(fn)
        path = EXPORTS / fn
        try:
            if fmt == "json":
                try:
                    d = json.loads(content) if isinstance(content, str) else content
                    path.write_text(json.dumps(d, indent=2), encoding="utf-8")
                except: path.write_text(content, encoding="utf-8")
            elif fmt == "csv":
                try:
                    d = json.loads(content) if isinstance(content, str) else content
                    if isinstance(d, list) and d:
                        with open(path, "w", newline="") as f:
                            w = csv.DictWriter(f, fieldnames=d[0].keys()); w.writeheader(); w.writerows(d)
                except: path.write_text(content, encoding="utf-8")
            elif fmt in ("yaml", "yml"):
                import yaml
                try:
                    d = json.loads(content) if isinstance(content, str) else content
                    with open(path, "w") as f: yaml.dump(d, f, default_flow_style=False)
                except: path.write_text(content, encoding="utf-8")
            elif fmt == "toml":
                import toml
                try:
                    d = json.loads(content) if isinstance(content, str) else content
                    with open(path, "w") as f: toml.dump(d, f)
                except: path.write_text(content, encoding="utf-8")
            elif fmt == "xml":
                xml_out = '<?xml version="1.0" encoding="UTF-8"?>\n<export><![CDATA[' + content + ']]></export>'
                path.write_text(xml_out, encoding="utf-8")
            elif fmt == "html":
                if not content.strip().startswith("<"):
                    content = f"<html><head><meta charset='utf-8'><title>{name or 'Export'}</title><style>body{{font-family:Arial;max-width:800px;margin:40px auto}}pre{{background:#f5f5f5;padding:15px}}</style></head><body><pre>{content}</pre></body></html>"
                path.write_text(content, encoding="utf-8")
            elif fmt == "ini":
                import configparser
                c = configparser.ConfigParser()
                try:
                    d = json.loads(content) if isinstance(content, str) else content
                    if isinstance(d, dict):
                        for s, v in d.items():
                            c[s] = v if isinstance(v, dict) else {"value": str(v)}
                        with open(path, "w") as f: c.write(f)
                except: path.write_text(content, encoding="utf-8")
            elif fmt == "rtf":
                rtf_text = '{\\rtf1\\ansi\\deff0 {\\fonttbl {\\f0 Arial;}}\n\\f0\\fs24 '
                rtf_text += content.replace(chr(10), '\\par ') + '}'
                path.write_text(rtf_text, encoding="utf-8")
            elif fmt == "svg":
                svg = '<svg xmlns="http://www.w3.org/2000/svg" width="800" height="600"><rect width="100%" height="100%" fill="white"/>'
                y = 20
                for line in content.split(chr(10)):
                    svg += '<text x="20" y="' + str(y) + '" font-family="Arial" font-size="14" fill="black">' + line.replace("&","&amp;").replace("<","&lt;") + '</text>'
                    y += 20
                svg += '</svg>'
                path.write_text(svg, encoding="utf-8")
            elif fmt == "pdf":
                from fpdf import FPDF
                pdf = FPDF(); pdf.add_page(); pdf.set_font("Arial", size=12)
                for line in content.split(chr(10)): pdf.cell(200, 10, txt=line, ln=True)
                pdf.output(str(path))
            elif fmt == "docx":
                from docx import Document
                doc = Document(); doc.add_heading(name or "Export", 0)
                for line in content.split(chr(10)):
                    if line.strip(): doc.add_paragraph(line)
                doc.save(str(path))
            elif fmt == "xlsx":
                import pandas as pd
                try:
                    d = json.loads(content) if isinstance(content, str) else content
                    if isinstance(d, list): pd.DataFrame(d).to_excel(path, index=False)
                    elif isinstance(d, dict): pd.DataFrame([d]).to_excel(path, index=False)
                except: path.write_text(content, encoding="utf-8")
            elif fmt == "pptx":
                from pptx import Presentation
                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = name or "Export"
                slide.placeholders[1].text = content
                prs.save(str(path))
            elif fmt in ("png", "jpg", "jpeg", "bmp", "gif", "tiff", "webp"):
                from PIL import Image, ImageDraw
                img = Image.new("RGB", (800, 600), color="white")
                draw = ImageDraw.Draw(img); y = 20
                for line in content.split(chr(10)): draw.text((20, y), line, fill="black"); y += 20
                img.save(str(path))
            elif fmt == "zip":
                tp = EXPORTS / f"{name or 'export'}_{ts}.txt"
                tp.write_text(content, encoding="utf-8")
                with zipfile.ZipFile(path, "w") as zf: zf.write(tp, tp.name)
                tp.unlink()
            elif fmt == "md":
                path.write_text(content, encoding="utf-8")
            else:
                path.write_text(content, encoding="utf-8")
            return f"Exported: {path.name} ({self._format_size(path.stat().st_size)})"
        except Exception as e:
            return f"Export failed: {e}"

    def handle(self, task):
        self.track_interaction()
        task = task.strip()
        if task.startswith("/"):
            parts = task.split(maxsplit=1)
            cmd = parts[0].lower()
            rest = parts[1] if len(parts) > 1 else ""
        else:
            cmd, rest = "", task
        try:
            if cmd == "/import" and rest: result = self._import(rest)
            elif cmd == "/export" and rest:
                fp = rest.split(maxsplit=1)
                result = self._export(fp[0], fp[1]) if len(fp) == 2 else "Usage: /export <format> <content>"
            elif cmd == "/convert" and rest:
                cp = rest.rsplit(maxsplit=1)
                if len(cp) == 2:
                    c = self._import(cp[0])
                    result = self._export(cp[1], c, Path(cp[0]).stem)
                else: result = "Usage: /convert <source> <target_format>"
            elif cmd == "/help":
                result = "FileClaw - Universal File Handler\n\n  IMPORT:  PDF, DOCX, XLSX, PPTX, RTF, images, archives, media + " + str(len(self.text_formats)) + " text formats\n  EXPORT:  PDF, DOCX, XLSX, PPTX, PNG/JPG/BMP/GIF/TIFF/WEBP, SVG, ZIP, JSON, CSV, YAML, TOML, XML, INI, RTF, MD, HTML, TXT\n  CONVERT: Any import format to any export format"
            elif cmd == "/stats":
                total = len(self.binary_importers) + len(self.text_formats)
                result = f"FileClaw | {total} formats | Import: {len(self.binary_importers)} binary + {len(self.text_formats)} text | Interactions: {self.state.get('interactions', 0)}"
            else:
                result = "FileClaw: Use /import, /export, /convert, or /help"
            return {"status": "success", "result": str(result)}
        except Exception as e:
            return {"status": "error", "result": str(e)}

_agent = FileClawAgent()
def process_task(task, agent=None):
    return _agent.handle(task)
